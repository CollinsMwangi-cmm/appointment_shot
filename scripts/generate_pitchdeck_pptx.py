# Script to generate a PPTX pitch deck from the markdown file using python-pptx
# Usage: python3 -m venv .venv && source .venv/bin/activate && pip install python-pptx markdown
#        python scripts/generate_pitchdeck_pptx.py docs/AppointmentShot_PitchDeck.md docs/AppointmentShot_PitchDeck.pptx

import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
import markdown

SLIDE_TITLE_STYLE = {
    'font_size': Pt(32),
}
SLIDE_BODY_STYLE = {
    'font_size': Pt(18),
}

def md_to_slides(md_text):
    # Very simple parser: splits on '---' for slides, first line is title if starts with 'Slide' or '#'
    slides = []
    parts = md_text.split('\n---\n')
    for part in parts:
        title = ''
        body = part.strip()
        lines = body.splitlines()
        if lines:
            # find first non-empty line as title if starts with 'Slide' or '#'
            first = lines[0].strip()
            if first.startswith('Slide') or first.startswith('#'):
                title = first
                body = '\n'.join(lines[1:]).strip()
        slides.append((title, body))
    return slides


def add_slide(prs, title, body):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    if title:
        slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for paragraph in body.split('\n'):
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        p.text = paragraph
        p.font.size = SLIDE_BODY_STYLE['font_size']


def main():
    if len(sys.argv) < 3:
        print('Usage: python scripts/generate_pitchdeck_pptx.py <input_md> <output_pptx>')
        sys.exit(1)
    md_in = sys.argv[1]
    pptx_out = sys.argv[2]
    with open(md_in, 'r', encoding='utf-8') as f:
        md_text = f.read()
    slides = md_to_slides(md_text)
    prs = Presentation()
    # remove default slide
    while len(prs.slides._sldIdLst):
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    for title, body in slides:
        add_slide(prs, title, body)
    prs.save(pptx_out)
    print(f'Generated {pptx_out}')

if __name__ == '__main__':
    main()